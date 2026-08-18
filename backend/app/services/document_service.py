"""文档入库与管理服务"""
import csv
import io
import re
import shutil
from pathlib import Path
from typing import Dict, List, Optional
from uuid import uuid4

from fastapi import BackgroundTasks, UploadFile
from sqlalchemy.orm import Session

from app.config import settings
from app.storage.sqlite_metadata import (
    DocumentChunkRecord,
    DocumentRecord,
    SessionLocal,
    init_metadata_db,
)
from app.storage.milvus_store import get_milvus_client
from app.utils.chunking import split_document_text
from app.utils.exceptions import DocumentException
from app.utils.logger import setup_logger

logger = setup_logger(__name__)


def _log_upload_step(document_id: str, step: str, message: str) -> None:
    """统一输出文档上传链路的步骤日志。"""
    logger.info(f"[文档上传][{step}] {message} (document_id={document_id})")


def _log_delete_step(document_id: str, step: str, message: str) -> None:
    """统一输出文档删除链路的步骤日志。"""
    logger.info(f"[文档删除][{step}] {message} (document_id={document_id})")


def _normalize_allowed_types() -> set[str]:
    """把配置里的允许文件类型清洗成可比较的集合。"""
    return {item.strip().lower().strip('"') for item in settings.ALLOWED_FILE_TYPES}


def _save_uploaded_file(file: UploadFile, document_id: str, original_name: str) -> Path:
    """保存上传文件到本地目录，并返回落盘后的路径。"""
    upload_dir = Path(settings.UPLOAD_DIR)
    upload_dir.mkdir(parents=True, exist_ok=True)

    stored_filename = f"{document_id}{Path(original_name).suffix.lower()}"
    stored_path = upload_dir / stored_filename

    file.file.seek(0)
    with stored_path.open("wb") as buffer:
        shutil.copyfileobj(file.file, buffer)

    _log_upload_step(
        document_id,
        "1/3",
        f"本地文件保存成功: {original_name} -> {stored_path}",
    )
    return stored_path


def _get_file_size(file_path: Path) -> int:
    """获取文件大小，便于做上传限制校验。"""
    return file_path.stat().st_size


def _delete_local_uploaded_file(file_path: Path, document_id: str, reason: str) -> bool:
    """删除上传后的临时原始文件，避免服务端长期保存用户文档。"""
    try:
        existed = file_path.exists()
        file_path.unlink(missing_ok=True)
        logger.info(
            "本地上传原始文件清理完成: deleted={} reason={} path={} (document_id={})",
            existed,
            reason,
            file_path,
            document_id,
        )
        return existed
    except Exception as exc:
        logger.warning(
            "本地上传原始文件清理失败: reason={} path={} error={} (document_id={})",
            reason,
            file_path,
            exc,
            document_id,
        )
        return False


def _extract_text_from_txt(file_path: Path) -> str:
    """从 TXT/MD 文件中提取文本，按常见编码做容错读取。"""
    for encoding in ("utf-8", "utf-8-sig", "gbk", "gb18030"):
        try:
            return file_path.read_text(encoding=encoding)
        except UnicodeDecodeError:
            continue
    raise DocumentException("TXT 文件编码无法识别")


def _extract_text_from_pdf(file_path: Path) -> str:
    """从 PDF 文件中抽取每一页的文本。"""
    try:
        from pypdf import PdfReader

        reader = PdfReader(str(file_path))
        pages = []
        for page in reader.pages:
            pages.append(page.extract_text() or "")
        return "\n".join(pages)
    except Exception as e:
        logger.warning(f"pypdf 解析失败，尝试 pypdfium2 兜底: {e}")

    try:
        import pypdfium2 as pdfium

        pdf = pdfium.PdfDocument(str(file_path))
        pages = []
        for page in pdf:
            textpage = page.get_textpage()
            pages.append(textpage.get_text_range() or "")
            textpage.close()
            page.close()
        pdf.close()
        return "\n".join(pages)
    except Exception as e:
        logger.warning(f"pypdfium2 解析失败，尝试 unstructured 兜底: {e}")
        return _extract_text_with_unstructured(file_path)


def _extract_text_from_docx(file_path: Path) -> str:
    """按原始顺序抽取 DOCX 段落和表格文本。"""
    from docx.oxml.table import CT_Tbl
    from docx.oxml.text.paragraph import CT_P
    from docx import Document as DocxDocument
    from docx.table import Table
    from docx.text.paragraph import Paragraph

    doc = DocxDocument(str(file_path))
    blocks: List[str] = []

    for child in doc.element.body.iterchildren():
        if isinstance(child, CT_P):
            paragraph = Paragraph(child, doc)
            text = " ".join(paragraph.text.split()).strip()
            if text:
                blocks.append(text)
            continue

        if not isinstance(child, CT_Tbl):
            continue

        table = Table(child, doc)
        rows: List[str] = []
        for row in table.rows:
            cells = [" ".join(cell.text.split()).strip() for cell in row.cells]
            if any(cells):
                rows.append(" | ".join(cells))
        if rows:
            blocks.append("[表格]\n" + "\n".join(rows))

    return "\n\n".join(blocks)


def _normalize_excel_cell(value, pandas_module) -> str:
    """把 Excel 单元格转换成稳定的单行文本。"""
    if value is None or pandas_module.isna(value):
        return ""
    if isinstance(value, float) and value.is_integer():
        return str(int(value))
    return " ".join(str(value).replace("\r", " ").replace("\n", " ").split())


def _serialize_excel_row(values: List[str]) -> str:
    """按 CSV 规则序列化一行，避免单元格中的逗号破坏列结构。"""
    buffer = io.StringIO()
    csv.writer(buffer, lineterminator="").writerow(values)
    return buffer.getvalue()


def _looks_numeric_excel_value(value: str) -> bool:
    """判断单元格是否更像数据值而不是字段名。"""
    normalized = (value or "").strip().replace(",", "")
    if not normalized:
        return False
    return bool(
        re.fullmatch(
            r"[-+]?(?:\d+(?:\.\d+)?|\.\d+)(?:%|元|万元|万)?",
            normalized,
        )
    )


def _infer_excel_header_index(rows: List[List[str]]) -> int:
    """从工作表行中推断真实表头位置，不依赖固定行号或业务字段名。"""
    nonempty_rows = [
        (index, row)
        for index, row in enumerate(rows)
        if any(cell.strip() for cell in row)
    ]
    if not nonempty_rows:
        return 0
    if len(nonempty_rows) == 1:
        return nonempty_rows[0][0]

    max_width = max(sum(bool(cell.strip()) for cell in row) for _, row in nonempty_rows)
    candidates = []
    for index, row in nonempty_rows:
        values = [cell.strip() for cell in row if cell.strip()]
        width = len(values)
        if width < 2:
            continue

        unique_ratio = len({value.casefold() for value in values}) / width
        text_ratio = sum(
            not _looks_numeric_excel_value(value) for value in values
        ) / width
        following = rows[index + 1 : index + 4]
        similar_rows = sum(
            sum(bool(cell.strip()) for cell in candidate)
            >= max(2, int(width * 0.5))
            for candidate in following
            if any(cell.strip() for cell in candidate)
        )

        score = (
            width * 2.0
            + unique_ratio * 3.0
            + text_ratio * 3.0
            + similar_rows * 2.0
        )
        if width >= max(2, int(max_width * 0.5)):
            score += 4.0
        candidates.append((score, width, -index, index))

    if candidates:
        return max(candidates)[3]
    return nonempty_rows[0][0]


def _normalize_excel_headers(headers: List[str]) -> List[str]:
    """清理空表头、Unnamed 表头并保证字段名唯一。"""
    normalized_headers: List[str] = []
    used_counts: Dict[str, int] = {}

    for index, header in enumerate(headers):
        value = " ".join((header or "").split()).strip()
        if not value or value.casefold().startswith("unnamed"):
            value = f"列{index + 1}"

        used_counts[value] = used_counts.get(value, 0) + 1
        count = used_counts[value]
        normalized_headers.append(value if count == 1 else f"{value}_{count}")

    return normalized_headers


def _extract_text_from_xlsx(file_path: Path) -> str:
    """抽取 Excel 工作表并自动识别真实表头，输出可稳定切块的结构化文本。"""
    import pandas as pd

    sheets = pd.read_excel(
        str(file_path),
        sheet_name=None,
        header=None,
        dtype=object,
        keep_default_na=False,
    )
    blocks = []

    for sheet_name, dataframe in sheets.items():
        raw_rows = [
            [
                _normalize_excel_cell(value, pd)
                for value in row
            ]
            for row in dataframe.itertuples(index=False, name=None)
        ]
        raw_rows = [
            row
            for row in raw_rows
            if any(cell.strip() for cell in row)
        ]
        if not raw_rows:
            continue

        max_width = max(len(row) for row in raw_rows)
        rows = [row + [""] * (max_width - len(row)) for row in raw_rows]
        while rows and all(not row[-1].strip() for row in rows):
            for row in rows:
                row.pop()

        header_index = _infer_excel_header_index(rows)
        headers = _normalize_excel_headers(rows[header_index])
        width = len(headers)

        blocks.append(f"[Sheet] {sheet_name}")
        blocks.append(_serialize_excel_row(headers))

        # 标题、说明和空行仍保留，但放在规范表头之后，避免被误识别成表头。
        for row in rows[:header_index]:
            values = [cell for cell in row[:width] if cell.strip()]
            if values:
                blocks.append(f"[说明] {' | '.join(values)}")

        for row in rows[header_index + 1 :]:
            normalized_row = row[:width] + [""] * max(0, width - len(row))
            if any(cell.strip() for cell in normalized_row):
                blocks.append(_serialize_excel_row(normalized_row))

    return "\n".join(blocks)


def _extract_text_from_pptx(file_path: Path) -> str:
    """从 PPTX 文件中抽取每页幻灯片文本。"""
    from pptx import Presentation

    presentation = Presentation(str(file_path))
    blocks = []
    for slide in presentation.slides:
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_text.append(shape.text.strip())
        if slide_text:
            blocks.append("\n".join(slide_text))
    return "\n\n".join(blocks)


def _extract_text_with_unstructured(file_path: Path) -> str:
    """用 unstructured 作为兜底解析器，覆盖更多文件格式。"""
    try:
        from unstructured.partition.auto import partition

        elements = partition(filename=str(file_path))
        return "\n".join(str(element).strip() for element in elements if str(element).strip())
    except Exception as e:
        raise DocumentException(f"文档解析失败: {e}")


def extract_text_from_file(file_path: Path) -> str:
    """根据文件后缀选择对应解析器，提取可向量化的文本内容。"""
    suffix = file_path.suffix.lower().lstrip(".")

    try:
        if suffix in {"txt", "md"}:
            return _extract_text_from_txt(file_path)
        if suffix == "pdf":
            return _extract_text_from_pdf(file_path)
        if suffix == "docx":
            return _extract_text_from_docx(file_path)
        if suffix in {"xlsx", "xls"}:
            return _extract_text_from_xlsx(file_path)
        if suffix == "pptx":
            return _extract_text_from_pptx(file_path)
        return _extract_text_with_unstructured(file_path)
    except DocumentException:
        raise
    except Exception as e:
        logger.error(f"提取文本失败: {str(e)}")
        raise DocumentException(f"提取文本失败: {str(e)}")


def _serialize_record(record: DocumentRecord) -> Dict:
    """把 ORM 记录转成接口可直接返回的字典。"""
    return {
        "document_id": record.document_id,
        "original_filename": record.original_filename,
        "stored_filename": record.stored_filename,
        "file_path": record.file_path,
        "content_type": record.content_type,
        "file_type": record.file_type,
        "file_size": record.file_size,
        "status": record.status,
        "chunk_count": record.chunk_count,
        "error_message": record.error_message,
        "created_at": record.created_at.isoformat() if record.created_at else None,
        "updated_at": record.updated_at.isoformat() if record.updated_at else None,
    }


def _update_document_record(document_id: str, **fields) -> None:
    """更新文档元数据中的状态字段，后台任务里复用。"""
    db: Session = SessionLocal()
    try:
        record = db.query(DocumentRecord).filter(DocumentRecord.document_id == document_id).first()
        if not record:
            return

        for key, value in fields.items():
            setattr(record, key, value)
        db.commit()
    finally:
        db.close()


def _replace_document_chunk_records(db: Session, document_id: str, chunks: List[str]) -> None:
    """用当前解析结果覆盖 SQLite 中的文档正文切片。"""
    db.query(DocumentChunkRecord).filter(DocumentChunkRecord.document_id == document_id).delete()
    db.add_all(
        DocumentChunkRecord(
            document_id=document_id,
            chunk_index=index,
            chunk_text=chunk,
        )
        for index, chunk in enumerate(chunks)
    )


def _get_document_content_from_sqlite(db: Session, document_id: str) -> str:
    """从 SQLite 切片表重组文档正文。"""
    chunks = (
        db.query(DocumentChunkRecord)
        .filter(DocumentChunkRecord.document_id == document_id)
        .order_by(DocumentChunkRecord.chunk_index.asc())
        .all()
    )
    return "\n\n".join(chunk.chunk_text for chunk in chunks if chunk.chunk_text.strip())


def _get_document_content_from_milvus(document_id: str) -> str:
    """兼容旧数据：从 Milvus chunk_text 重组正文。"""
    try:
        client = get_milvus_client()
        rows = client.query(
            collection_name=settings.MILVUS_DOC_COLLECTION_NAME,
            filter=f'document_id == "{document_id}"',
            output_fields=["chunk_index", "chunk_text"],
            limit=10000,
        )
    except Exception as exc:
        logger.warning("从 Milvus 读取文档正文切片失败，已跳过: {} (document_id={})", exc, document_id)
        return ""

    sorted_rows = sorted(rows or [], key=lambda item: int(item.get("chunk_index") or 0))
    return "\n\n".join(
        str(row.get("chunk_text") or "").strip()
        for row in sorted_rows
        if str(row.get("chunk_text") or "").strip()
    )


def get_document_content(document_id: str) -> Dict:
    """获取文档正文预览，优先使用 SQLite 切片，避免依赖本地原始文件。"""
    init_metadata_db()
    db: Session = SessionLocal()
    try:
        record = db.query(DocumentRecord).filter(DocumentRecord.document_id == document_id).first()
        if not record:
            raise DocumentException("文档不存在")

        content = _get_document_content_from_sqlite(db, document_id)
        if not content:
            content = _get_document_content_from_milvus(document_id)

        if not content:
            file_path = Path(record.file_path)
            if file_path.exists():
                content = extract_text_from_file(file_path)

        if not content:
            raise DocumentException("文档正文不存在或尚未处理完成")

        return {
            "document_id": record.document_id,
            "original_filename": record.original_filename,
            "file_type": record.file_type,
            "status": record.status,
            "content": content,
            "content_type": record.content_type,
            "updated_at": record.updated_at,
        }
    finally:
        db.close()


def _insert_document_chunks_to_milvus(
    document_id: str,
    chunks: List[str],
    source_name: str,
    file_type: str,
    content_type: Optional[str],
) -> List[str]:
    """使用 MilvusClient 原生写入文档切块。"""
    from app.core.embeddings import get_default_embeddings

    logger.info(f"开始初始化 BGE Embedding (document_id={document_id}, chunks={len(chunks)})")
    embeddings = get_default_embeddings()
    logger.info(f"BGE Embedding 就绪，开始向量化 (document_id={document_id}, chunks={len(chunks)})")
    vectors = embeddings.embed_documents(chunks)
    logger.info(f"BGE 向量化完成 (document_id={document_id}, vectors={len(vectors)})")
    client = get_milvus_client()

    rows = []
    for index, chunk in enumerate(chunks):
        rows.append(
            {
                "text": chunk,
                "vector": vectors[index],
                "document_id": document_id,
                "chunk_index": index,
                "source_name": source_name,
                "chunk_text": chunk,
                "file_type": file_type,
                "content_type": content_type or "",
            }
        )

    from app.storage.milvus_store import insert_rows_with_retry

    ids = insert_rows_with_retry(
        client=client,
        collection_name=settings.MILVUS_DOC_COLLECTION_NAME,
        rows=rows,
    )
    client.flush(collection_name=settings.MILVUS_DOC_COLLECTION_NAME)
    logger.info(f"Milvus 文档切块写入成功 (document_id={document_id}, chunks={len(ids)})")
    return ids


def _delete_document_chunks_from_milvus(
    document_id: str,
) -> bool:
    """使用 MilvusClient 原生删除文档切块。"""

    collection_name = settings.MILVUS_DOC_COLLECTION_NAME
    client = get_milvus_client()
    if not client.has_collection(collection_name):
        logger.warning(f"[文档删除][2/4] Milvus collection 不存在，无需删除切块 ({collection_name}, document_id={document_id})")
        return False

    client.delete(
        collection_name=collection_name,
        filter=f'document_id == "{document_id}"',
    )
    return True


def _document_delete_requested(document_id: str) -> bool:
    """后台任务检查文档是否已经被用户删除。"""
    db: Session = SessionLocal()
    try:
        record = db.query(DocumentRecord).filter(DocumentRecord.document_id == document_id).first()
        return record is None or record.status in {"deleting", "deleted"}
    finally:
        db.close()


def _validate_document_delete_prerequisites(db: Session, document_id: str) -> tuple[DocumentRecord, Path]:
    """删除前校验 SQLite 与本地文件是否存在。"""

    record = db.query(DocumentRecord).filter(DocumentRecord.document_id == document_id).first()
    if not record:
        raise DocumentException("文档不存在")
    if record.status in {"deleting", "deleted"}:
        raise DocumentException(f"文档当前状态为 {record.status}，无法重复删除")

    return record, Path(record.file_path)


def _process_document_upload(
    document_id: str,
    stored_path: Path,
    original_name: str,
    file_type: str,
    content_type: Optional[str],
) -> None:
    """后台执行文档解析、切块与向量写入。"""
    db: Session = SessionLocal()
    try:
        logger.info(
            f"后台处理开始 (document_id={document_id}, file={original_name}, path={stored_path})"
        )

        logger.info(f"开始抽取文本 (document_id={document_id})")
        text = extract_text_from_file(stored_path)
        logger.info(
            f"文本抽取完成 (document_id={document_id}, text_len={len(text)})"
        )
        if not text.strip():
            raise DocumentException("文档内容为空，无法生成向量")

        logger.info(f"开始切块 (document_id={document_id})")
        chunks = split_document_text(text, file_type=file_type)
        logger.info(
            f"切块完成 (document_id={document_id}, chunk_count={len(chunks)})"
        )
        if not chunks:
            raise DocumentException("文档切块失败，未生成有效内容")

        if _document_delete_requested(document_id):
            logger.info(f"文档已请求删除，跳过后台向量写入 (document_id={document_id})")
            return

        _log_upload_step(
            document_id,
            "3/3",
            f"开始写入 Milvus (chunks={len(chunks)})",
        )
        _insert_document_chunks_to_milvus(
            document_id=document_id,
            chunks=chunks,
            source_name=original_name,
            file_type=file_type,
            content_type=content_type,
        )
        _log_upload_step(
            document_id,
            "3/3",
            f"Milvus 写入成功: chunks={len(chunks)}",
        )

        if _document_delete_requested(document_id):
            _delete_document_chunks_from_milvus(document_id)
            logger.info(f"文档已在向量写入期间删除，已清理向量块 (document_id={document_id})")
            return

        logger.info(f"开始更新文档状态为 ready (document_id={document_id})")
        record = db.query(DocumentRecord).filter(DocumentRecord.document_id == document_id).first()
        if record:
            _replace_document_chunk_records(db, document_id, chunks)
            record.status = "ready"
            record.chunk_count = len(chunks)
            record.error_message = None
            db.commit()
            logger.info(f"文档状态更新完成 (document_id={document_id}, status=ready, chunks={len(chunks)})")
        logger.info(f"✓ 文档后台处理完成 (document_id={document_id}, chunks={len(chunks)})")
    except Exception as e:
        db.rollback()
        try:
            logger.info(f"开始更新文档状态为 failed (document_id={document_id})")
            record = db.query(DocumentRecord).filter(DocumentRecord.document_id == document_id).first()
            if record:
                record.status = "failed"
                record.error_message = str(e)
                db.commit()
                logger.info(f"文档状态更新完成 (document_id={document_id}, status=failed)")
        except Exception:
            db.rollback()
        logger.exception(f"文档后台处理失败 (document_id={document_id}): {e}")
    finally:
        db.close()
        _delete_local_uploaded_file(stored_path, document_id, "background_processed")


def ingest_document(file: UploadFile, background_tasks: Optional[BackgroundTasks] = None) -> Dict:
    """上传文档，解析文本，切块后写入 SQLite 元数据和 Milvus 向量库。"""
    if not file.filename:
        raise DocumentException("文件名不能为空")

    # 1) 先做后缀校验，避免不支持的格式进入后续链路
    original_name = Path(file.filename).name
    file_type = Path(original_name).suffix.lower().lstrip(".")
    allowed_types = _normalize_allowed_types()
    if file_type not in allowed_types:
        raise DocumentException(f"不支持的文件类型: {file_type}")

    # 2) 给文件生成内部 ID，保证系统内文件名唯一且可追踪
    document_id = uuid4().hex
    _log_upload_step(document_id, "1/3", f"开始接收上传文件: {original_name}")
    stored_path = _save_uploaded_file(file, document_id, original_name)
    file_size = _get_file_size(stored_path)

    # 3) 先按字节大小拦截超限文件，避免后续解析和向量化浪费资源
    if file_size > settings.MAX_UPLOAD_SIZE:
        stored_path.unlink(missing_ok=True)
        raise DocumentException("文件超过最大上传限制")

    # 4) 确保 SQLite 元数据表已创建，然后写入一条 processing 记录
    init_metadata_db()
    db: Session = SessionLocal()

    record = DocumentRecord(
        document_id=document_id,
        original_filename=original_name,
        stored_filename=stored_path.name,
        file_path=str(stored_path.resolve()),
        content_type=file.content_type,
        file_type=file_type,
        file_size=file_size,
        status="processing",
        chunk_count=0,
    )

    try:
        db.add(record)
        db.commit()
        _log_upload_step(
            document_id,
            "2/3",
            f"SQLite 元数据写入成功: status=processing, file_size={file_size}",
        )

        if background_tasks is not None:
            background_tasks.add_task(
                _process_document_upload,
                document_id=document_id,
                stored_path=stored_path,
                original_name=original_name,
                file_type=file_type,
                content_type=file.content_type,
            )
            db.refresh(record)
            logger.info(f"[文档上传] 后台处理任务已提交 (document_id={document_id})")
            return _serialize_record(record)

        text = extract_text_from_file(stored_path)
        if not text.strip():
            raise DocumentException("文档内容为空，无法生成向量")

        chunks = split_document_text(text, file_type=file_type)
        if not chunks:
            raise DocumentException("文档切块失败，未生成有效内容")

        _insert_document_chunks_to_milvus(
            document_id=document_id,
            chunks=chunks,
            source_name=original_name,
            file_type=file_type,
            content_type=file.content_type,
        )

        record.status = "ready"
        record.chunk_count = len(chunks)
        record.error_message = None
        _replace_document_chunk_records(db, document_id, chunks)
        db.commit()
        db.refresh(record)
        _log_upload_step(
            document_id,
            "3/3",
            f"Milvus 写入成功并更新文档状态: ready, chunk_count={len(chunks)}",
        )

        return _serialize_record(record)
    except Exception as e:
        db.rollback()
        try:
            record.status = "failed"
            record.error_message = str(e)
            db.add(record)
            db.commit()
        except Exception:
            db.rollback()
        logger.exception(f"文档入库失败 (document_id={document_id}): {e}")
        raise
    finally:
        db.close()
        if "stored_path" in locals() and background_tasks is None:
            _delete_local_uploaded_file(stored_path, document_id, "sync_processed")


def list_documents(skip: int = 0, limit: int = 10) -> Dict:
    """从 SQLite 元数据返回完整文档生命周期状态。"""
    init_metadata_db()
    resolved_skip = max(skip, 0)
    resolved_limit = max(limit, 0)
    db: Session = SessionLocal()
    try:
        query = db.query(DocumentRecord).order_by(DocumentRecord.created_at.desc())
        total = query.count()
        records = query.offset(resolved_skip).limit(resolved_limit).all()
        return {
            "documents": [_serialize_record(record) for record in records],
            "total": total,
            "skip": resolved_skip,
            "limit": resolved_limit,
        }
    finally:
        db.close()


def request_document_deletion(document_id: str) -> Dict:
    """先把文档标记为 deleting，供接口快速返回并提交后台清理。"""
    init_metadata_db()
    db: Session = SessionLocal()
    try:
        _log_delete_step(document_id, "1/4", "收到删除请求")
        record, file_path = _validate_document_delete_prerequisites(db, document_id)
        _log_delete_step(
            document_id,
            "1/4",
            f"SQLite 元数据校验通过: status={record.status}, file_path={file_path}",
        )

        record.status = "deleting"
        record.error_message = None
        db.commit()
        _log_delete_step(document_id, "1/4", "SQLite 状态已更新为 deleting")
        payload = _serialize_record(record)
        payload["milvus_deleted"] = False
        payload["file_deleted"] = False
        return payload
    except Exception:
        db.rollback()
        raise
    finally:
        db.close()


def delete_document(document_id: str, already_marked: bool = False) -> Dict:
    """后台删除 SQLite 元数据、Milvus 向量块和本地原始文件。"""
    init_metadata_db()
    db: Session = SessionLocal()
    try:
        _log_delete_step(document_id, "1/4", "开始执行后台删除")
        record = db.query(DocumentRecord).filter(DocumentRecord.document_id == document_id).first()
        if not record:
            raise DocumentException("文档不存在")
        if record.status == "deleted":
            raise DocumentException("文档已经删除")
        if record.status != "deleting":
            if already_marked:
                raise DocumentException(f"文档当前状态为 {record.status}，无法执行后台删除")
            record.status = "deleting"
            record.error_message = None
            db.commit()
            _log_delete_step(document_id, "1/4", "SQLite 状态已更新为 deleting")

        file_path = Path(record.file_path)

        _log_delete_step(document_id, "2/4", "开始删除 Milvus 文档切块")
        try:
            milvus_deleted = _delete_document_chunks_from_milvus(document_id)
            _log_delete_step(
                document_id,
                "2/4",
                f"Milvus 文档切块删除完成: deleted={milvus_deleted}",
            )
        except Exception as exc:
            logger.exception(f"[文档删除][2/4] Milvus 文档切块删除失败 (document_id={document_id}): {exc}")
            raise

        _log_delete_step(document_id, "3/4", f"开始删除本地文件: {file_path}")
        try:
            file_deleted = file_path.exists()
            file_path.unlink(missing_ok=True)
            _log_delete_step(
                document_id,
                "3/4",
                f"本地文件删除完成: deleted={file_deleted}",
            )
        except Exception as exc:
            logger.exception(f"[文档删除][3/4] 本地文件删除失败 (document_id={document_id}): {exc}")
            raise

        _log_delete_step(document_id, "4/4", "开始删除 SQLite 文档记录")
        try:
            payload = _serialize_record(record)
            payload["status"] = "deleted"
            db.query(DocumentChunkRecord).filter(DocumentChunkRecord.document_id == document_id).delete()
            db.delete(record)
            db.commit()
            _log_delete_step(document_id, "4/4", "SQLite 文档记录删除完成")
        except Exception as exc:
            logger.exception(f"[文档删除][4/4] SQLite 文档记录删除失败 (document_id={document_id}): {exc}")
            raise
        payload["milvus_deleted"] = milvus_deleted
        payload["file_deleted"] = file_deleted

        logger.info(f"✓ 文档删除完成 (document_id={document_id})")
        return payload
    except Exception as e:
        db.rollback()
        try:
            record = db.query(DocumentRecord).filter(DocumentRecord.document_id == document_id).first()
            if record:
                record.status = "ready"
                record.error_message = f"删除失败: {str(e)}"
                db.commit()
                _log_delete_step(document_id, "1/4", "后台删除失败，文档状态已恢复为 ready")
        except Exception:
            db.rollback()
        logger.exception(f"文档删除失败 (document_id={document_id}): {e}")
        raise
    finally:
        db.close()
